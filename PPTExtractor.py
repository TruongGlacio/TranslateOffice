from pptx import Presentation
import os
import GeneralFuntions

generalFunctions=GeneralFuntions.General()
class PPTExtractor:
    def __init__(self):
        pass 
    
    def ReadTextFromPPT(self,inputPath, outputTextFilePath):
        textList=list()
        textList.append(PPTExtractor.ReadTextTable(self,inputPath))
        textList.append(PPTExtractor.ReadTextNormal(self,inputPath))
        textListJson=generalFunctions.ConvertStringToJson(textList)
        generalFunctions.WriteTextFile(textListJson, outputTextFilePath)
        return textListJson
    def ReadTextNormal(self,inputPath):
        textList=list()
        prs = Presentation(inputPath)    
        print('function ReadTextNormal' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    textList.append(paragraph.text)
                    print('output text:',paragraph.text)
        return textList                
    
    def ReadTextTable(self,inputPath):
        textList=list()
        prs = Presentation(inputPath)    
        print('function ReadTextTable' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue  
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        textList.append(c.text_frame.text)
                        print('output text:',c.text_frame.text)
        return textList
                        
    def WriteTextFromJson(seft, inputPath,textListJson):
        textList=list()
        textTableList=list()
        textNomalList=list()
        textList=generalFunctions.ConvertJsonToString(textListJson)
        textTableList=textList[0]
        textNomalList=textList[1]
        print('textNomalList:', textNomalList)
        print('textTableList:',textTableList)
        PPTExtractor.WriteTextTable(seft, inputPath,textTableList)
        PPTExtractor.WriteTextNormal(seft, inputPath,textNomalList)     
    def WriteTextNormal(self,inputPath, inputStringList):
        i=0
        prs = Presentation(inputPath)
        print('function ReadTextNormal' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                        if(paragraph.text==''):
                            continue
                        else:
                            paragraph.text=inputStringList[i]
                        i=i+1    
                        print('output text:',paragraph.text)
        prs.save(inputPath)
        
    def WriteTextTable(self,inputPath, inputStringList):
        i=0
        prs = Presentation(inputPath)    
        print('function ReadTextTable' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue  
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        c.text_frame.text=inputStringList[i]
                        print('output text:',c.text_frame.text)
                        i=i+1 
        prs.save(inputPath)        
