from pptx import Presentation
import os
import GeneralFuntions

generalFunctions=GeneralFuntions.General()
LengthOfList=0
filePath=''
class PPTExtractor:
    def __init__(self):
        pass 
    
    def GetLength(self):
        return LengthOfList
    
    def SetLength(self, inputLength):
        global LengthOfList
        LengthOfList=inputLength
        
    def GetZipFolder(self): 
        return filePath
    
    def SetZipFolder(self, inpuzipfolder):
        global filePath
        filePath=inpuzipfolder
        
    def ReadTextFromPPT(self,inputPath, outputTextFilePath):
        textList=list()
        textListTable=list()
        textListNormal=list()
        PPTExtractor.SetZipFolder(self,inputPath)
        textListTable=PPTExtractor.ReadTextTable(self,inputPath)
        textListNormal=PPTExtractor.ReadTextNormal(self,inputPath)
        for text in textListTable:
            textList.append(text)
        for text in textListNormal:
            textList.append(text)
        PPTExtractor.SetLength(self, len(textList))
        textListJson=generalFunctions.ConvertStringToJson(textList)
        generalFunctions.WriteTextFile(textListJson, outputTextFilePath)
        return outputTextFilePath
    def ReadTextNormal(self,inputPath):
        textList=list()
        prs = Presentation(inputPath)    
        print('function ReadTextNormal' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if(paragraph.text==None):
                        continue
                    textList.append(paragraph.text)
                    print('output text:',paragraph.text)
        return textList                
    
    def ReadTextTable(self,inputPath):
        
        textList=list()
        prs = Presentation(inputPath)    
        print('function ReadTextTable' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue  
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        if(c.text_frame.text==None):
                            continue
                        textList.append(c.text_frame.text)
                        print('output text:',c.text_frame.text)
        return textList
    
    def WriteTextFromJson(seft, inputPath,textListJson,count):
        textList=list()
        textList=generalFunctions.ConvertJsonToString(textListJson)
        lengthString=PPTExtractor.GetLength(seft)
        print('textList:', textList)        
        if(count==len(textList)): 
            print('number item correct, go write data to file')            
            i= PPTExtractor.WriteTextTable(seft, inputPath,textList)
            PPTExtractor.WriteTextNormal(seft, inputPath,textList,i)    
            return inputPath
        else:
            print('number item incorrect, go exit')
            return None
    def WriteTextNormal(self,inputPath, inputStringList,i):
        prs = Presentation(inputPath)
        print('function ReadTextNormal' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                        if(paragraph.text==None):
                            continue
                        paragraph.text=inputStringList[i]
                        i=i+1    
                        print('output text:',paragraph.text)
        prs.save(inputPath)
        return i
        
    def WriteTextTable(self,inputPath, inputStringList):
        i=0
        prs = Presentation(inputPath)    
        print('function ReadTextTable' )
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue  
                table = shape.table
                for r in table.rows:
                    for c in r.cells:
                        if(c.text_frame.text==None):
                            continue
                        c.text_frame.text=inputStringList[i]
                        print('output text:',c.text_frame.text)
                        i=i+1 
        prs.save(inputPath)  
        return i
